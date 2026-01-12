from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Body
from rapidocr_onnxruntime import RapidOCR
from pdf2image import convert_from_bytes
from PIL import Image
import numpy as np
import io
import time
import logging
import requests
import os
import tempfile
import subprocess

# 新增库引入
from openpyxl import load_workbook
from pptx import Presentation
import xlrd  # 用于处理 .xls

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="OCR Service")

# 1. 全局初始化 OCR 模型
try:
    logger.info("正在加载OCR模型...")
    ocr_engine = RapidOCR()
    logger.info("OCR模型加载完成！")
except Exception as e:
    logger.error(f"模型加载失败: {e}")
    raise e

# --- 核心处理逻辑函数 ---

def ocr_single_image(img_np):
    """辅助函数：处理单张图片的 Numpy 数组并返回文本"""
    try:
        result, _ = ocr_engine(img_np)
        txt_res = []
        if result:
            for line in result:
                txt_res.append(line[1])
        return "\n".join(txt_res)
    except Exception as e:
        logger.error(f"单图识别出错: {e}")
        return ""

def extract_from_excel(file_bytes):
    """解析 Excel (.xlsx) 文本"""
    try:
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
        texts = []
        for sheet in wb.worksheets:
            texts.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"Excel(.xlsx)解析失败: {e}")
        raise HTTPException(status_code=400, detail=f"xlsx解析失败: {str(e)}")

def extract_from_xls(file_bytes):
    """解析旧版 Excel (.xls) 文本"""
    try:
        # xlrd 支持直接从内存读取
        wb = xlrd.open_workbook(file_contents=file_bytes)
        texts = []
        for sheet in wb.sheets():
            texts.append(f"--- Sheet: {sheet.name} ---")
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                # xlrd 的 cell 对象需要取 .value
                row_text = " ".join([str(c.value) for c in row if c.value not in ('', None)])
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except Exception as e:
        logger.error(f"Excel(.xls)解析失败: {e}")
        raise HTTPException(status_code=400, detail=f"xls解析失败: {str(e)}")

def extract_from_ppt(file_bytes):
    """解析 PPT (.pptx) 文本"""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for i, slide in enumerate(prs.slides):
            page_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    page_content.append(shape.text)
            if page_content:
                texts.append(f"--- Slide {i+1} ---")
                texts.append("\n".join(page_content))
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PPT(.pptx)解析失败: {e}")
        raise HTTPException(status_code=400, detail="pptx解析失败")

def extract_from_ppt_legacy(file_bytes):
    """解析旧版 PPT (.ppt) 文本"""
    # .ppt 是二进制 OLE 格式，纯 Python 库支持极差
    # 这里使用系统级工具 catppt (来自 catdoc 包) 进行提取
    temp_file = None
    try:
        # 创建临时文件，因为 catppt 需要文件路径
        with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as tmp:
            tmp.write(file_bytes)
            temp_file = tmp.name

        # 调用系统命令 catppt
        # -d utf-8 尝试强制输出 utf-8 (视 catdoc 版本而定，通常默认即可)
        process = subprocess.Popen(['catppt', temp_file], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        stdout, stderr = process.communicate()

        if process.returncode != 0:
            logger.error(f"catppt error: {stderr.decode()}")
            raise Exception("解析进程异常退出")

        # 尝试解码
        return stdout.decode('utf-8', errors='ignore')

    except Exception as e:
        logger.error(f"PPT(.ppt)解析失败: {e}")
        raise HTTPException(status_code=400, detail="ppt解析失败 (请检查文件是否损坏)")
    finally:
        # 清理临时文件
        if temp_file and os.path.exists(temp_file):
            os.remove(temp_file)

def process_file_bytes(file_bytes, filename):
    """通用文件流处理入口"""
    filename = filename.lower()

    # 1. PDF 处理
    if filename.endswith('.pdf'):
        try:
            images = convert_from_bytes(file_bytes, dpi=200, fmt='jpeg')
            texts = []
            for i, img in enumerate(images):
                img_np = np.array(img)
                texts.append(f"--- Page {i+1} ---\n" + ocr_single_image(img_np))
            return "\n\n".join(texts)
        except Exception as e:
            raise HTTPException(status_code=500, detail="PDF解析失败")

    # 2. Excel 处理
    elif filename.endswith('.xlsx'):
        return extract_from_excel(file_bytes)
    elif filename.endswith('.xls'):
        return extract_from_xls(file_bytes)

    # 3. PPT 处理
    elif filename.endswith('.pptx'):
        return extract_from_ppt(file_bytes)
    elif filename.endswith('.ppt'):
        return extract_from_ppt_legacy(file_bytes)

    # 4. 图片处理
    else:
        try:
            img = Image.open(io.BytesIO(file_bytes))
            img_np = np.array(img)
            return ocr_single_image(img_np)
        except Exception:
            raise HTTPException(status_code=400, detail=f"不支持的文件格式: {filename}")

# --- 接口定义 ---

@app.get("/health")
def health_check():
    return {"status": "ok"}

@app.post("/ocr")
async def recognize_file(file: UploadFile = File(...)):
    """
    文件上传接口
    支持: 图片, PDF, Excel(.xlsx/.xls), PPT(.pptx/.ppt)
    """
    start_time = time.time()
    try:
        content = await file.read()
        result_text = process_file_bytes(content, file.filename)

        cost = (time.time() - start_time) * 1000
        return {"code": 200, "data": result_text, "cost_time_ms": cost}
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"处理异常: {e}")
        return {"code": 500, "msg": str(e)}

@app.post("/ocr/url")
async def recognize_url(url: str = Body(..., embed=True)):
    """
    URL 识别接口
    Body参数: {"url": "http://example.com/file.png"}
    """
    start_time = time.time()
    try:
        logger.info(f"正在下载文件: {url}")
        resp = requests.get(url, timeout=15)
        if resp.status_code != 200:
            raise HTTPException(status_code=400, detail=f"下载失败，状态码: {resp.status_code}")

        filename = os.path.basename(url.split("?")[0])
        if not filename or "." not in filename:
            content_type = resp.headers.get('Content-Type', '')
            if 'pdf' in content_type: filename = 'temp.pdf'
            elif 'sheet' in content_type or 'excel' in content_type:
                # 简单判断，这里默认给 xlsx，如果是 xls 可能需要更复杂的 magic number 判断
                filename = 'temp.xlsx'
            elif 'presentation' in content_type: filename = 'temp.pptx'
            else: filename = 'temp.jpg'

        result_text = process_file_bytes(resp.content, filename)

        cost = (time.time() - start_time) * 1000
        return {"code": 200, "data": result_text, "cost_time_ms": cost}

    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"URL处理异常: {e}")
        return {"code": 500, "msg": str(e)}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=9000)